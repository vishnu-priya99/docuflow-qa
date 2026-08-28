"""PPTX parsing: one ParsedUnit per slide."""
from __future__ import annotations

import io

from pptx import Presentation

from app.services.ingestion.table_rendering import cell_text
from app.services.ingestion.units import ParsedUnit


def parse_pptx(content: bytes) -> list[ParsedUnit]:
    presentation = Presentation(io.BytesIO(content))
    units: list[ParsedUnit] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        title = None
        if slide.shapes.title is not None and slide.shapes.title.text.strip():
            title = slide.shapes.title.text.strip()

        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                texts.append(shape.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    # A genuinely blank cell renders as "(blank)" rather
                    # than being dropped - see table_rendering.cell_text.
                    raw_cells = [c.text.strip() for c in row.cells]
                    if any(raw_cells):
                        texts.append(" | ".join(cell_text(c) for c in raw_cells))

        combined = "\n".join(texts).strip()
        if combined:
            units.append(
                ParsedUnit(
                    text=combined,
                    section=title,
                    slide_number=slide_index,
                    slide_title=title,
                )
            )
    return units
