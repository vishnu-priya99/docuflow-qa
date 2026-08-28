"""Helpers that build small, real PDF/DOCX/PPTX/XLSX/CSV byte payloads for
tests, so ingestion is exercised against real parsers rather than stubs."""
from __future__ import annotations

import io

import pandas as pd
from docx import Document
from fpdf import FPDF
from pptx import Presentation


def make_pdf_bytes(pages: list[str]) -> bytes:
    """One page per string; the first line of each page becomes its heading
    if it looks like one (short, title/upper case)."""
    pdf = FPDF()
    for page_text in pages:
        pdf.add_page()
        pdf.set_font("Helvetica", size=12)
        pdf.multi_cell(0, 8, text=page_text)
    return bytes(pdf.output())


def make_pdf_with_no_text_bytes() -> bytes:
    """A page with no text layer at all - simulates a scanned/image-only PDF."""
    pdf = FPDF()
    pdf.add_page()
    return bytes(pdf.output())


def make_docx_bytes(blocks: list[tuple[str, str]]) -> bytes:
    """blocks: list of (style, text). style is "Heading 1" or "Normal"."""
    document = Document()
    for style, text in blocks:
        document.add_paragraph(text, style=style)
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def make_pptx_bytes(slides: list[tuple[str, str]]) -> bytes:
    """slides: list of (title, body_text)."""
    presentation = Presentation()
    layout = presentation.slide_layouts[1]  # title + content
    for title, body in slides:
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title
        body_placeholder = slide.placeholders[1]
        body_placeholder.text = body
    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


def make_txt_bytes(text: str) -> bytes:
    return text.encode("utf-8")


def make_csv_bytes(df: pd.DataFrame) -> bytes:
    buf = io.BytesIO()
    df.to_csv(buf, index=False)
    return buf.getvalue()


def make_xlsx_bytes(sheets: dict[str, pd.DataFrame]) -> bytes:
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        for name, df in sheets.items():
            df.to_excel(writer, sheet_name=name, index=False)
    return buf.getvalue()


def make_messy_xlsx_bytes(sheet_name: str, df: pd.DataFrame, *, title: str) -> bytes:
    """A realistic "business report" shape: a merged title banner, a blank
    row, THEN the real headers, plus a trailing footnote row - as opposed
    to make_xlsx_bytes' clean header-in-row-0 output."""
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet_name

    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=len(df.columns))
    ws.cell(row=1, column=1, value=title)
    # row 2 left blank
    header_row = 3
    for col_idx, col_name in enumerate(df.columns, start=1):
        ws.cell(row=header_row, column=col_idx, value=col_name)
    for row_offset, (_, row) in enumerate(df.iterrows()):
        for col_idx, value in enumerate(row, start=1):
            ws.cell(row=header_row + 1 + row_offset, column=col_idx, value=value)
    footnote_row = header_row + 1 + len(df)
    ws.cell(row=footnote_row + 1, column=1, value="Note: figures subject to revision")

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()
